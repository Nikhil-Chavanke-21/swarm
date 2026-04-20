from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0D, 0x0D)   # near-black
ACCENT  = RGBColor(0x7C, 0x3A, 0xED)   # purple
ACCENT2 = RGBColor(0x06, 0xB6, 0xD4)   # cyan
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xA0, 0xA0, 0xB0)
CARD    = RGBColor(0x1A, 0x1A, 0x2E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def txbox(slide, text, l, t, w, h, size=24, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def rect(slide, l, t, w, h, fill=CARD, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    return shp

def accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(1, 0, H - Inches(0.12), W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def slide_number(slide, n):
    txbox(slide, str(n), W - Inches(0.6), H - Inches(0.4), Inches(0.5), Inches(0.3),
          size=10, color=GREY, align=PP_ALIGN.RIGHT)

# ── Slide 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)

# Purple glow blob (rectangle, low-opacity simulation via layered rects)
blob = s.shapes.add_shape(9, Inches(3), Inches(1.5), Inches(7), Inches(5))
blob.fill.solid()
blob.fill.fore_color.rgb = RGBColor(0x3B, 0x00, 0x8A)
blob.line.fill.background()

txbox(s, "🐝", Inches(1.0), Inches(1.6), Inches(1.5), Inches(1.2), size=72, align=PP_ALIGN.CENTER)
txbox(s, "SWARM", Inches(2.3), Inches(1.5), Inches(9), Inches(1.8),
      size=90, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txbox(s, "Run many Claude Code agents in parallel — each in its own sandbox.",
      Inches(2.3), Inches(3.2), Inches(8.5), Inches(1.0),
      size=26, color=ACCENT2, align=PP_ALIGN.LEFT)
txbox(s, "Hackathon Demo  •  2026", Inches(2.3), Inches(4.4), Inches(8), Inches(0.5),
      size=16, color=GREY, align=PP_ALIGN.LEFT)
accent_bar(s)
slide_number(s, 1)

# ── Slide 2 — The Problem ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
accent_bar(s)
txbox(s, "😤  The Itch", Inches(0.6), Inches(0.4), Inches(10), Inches(0.9),
      size=40, bold=True, color=WHITE)

problems = [
    ("🖥️",  "Terminal chaos",         "Juggling dozens of tmux panes just to run a few agents"),
    ("🔒",  "1 session at a time",    "Each workspace is locked — no parallelism"),
    ("👀",  "Constant babysitting",   "You have to watch every agent or miss failures"),
    ("📋",  "No specialisation",      "Every agent re-explains itself — zero reuse"),
    ("🤝",  "No sharing",             "Great workflows die on your laptop"),
]

for i, (emoji, title, body) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.35)
    y = Inches(1.5 + row * 2.5)
    card = rect(s, x, y, Inches(4.1), Inches(2.2), fill=CARD, line=ACCENT)
    txbox(s, emoji, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.6), size=28)
    txbox(s, title, x + Inches(0.9), y + Inches(0.15), Inches(3.0), Inches(0.5),
          size=16, bold=True, color=WHITE)
    txbox(s, body, x + Inches(0.2), y + Inches(0.7), Inches(3.7), Inches(1.3),
          size=13, color=GREY)

slide_number(s, 2)

# ── Slide 3 — The Solution ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
accent_bar(s)
txbox(s, "💡  Swarm — The Solution", Inches(0.6), Inches(0.3), Inches(12), Inches(0.9),
      size=40, bold=True, color=WHITE)

pills = [
    (ACCENT,  "Define once",      "Specialised agent: instructions, tools, repos, MCPs"),
    (ACCENT2, "Run in parallel",  "Many instances, each with its own isolated sandbox"),
    (RGBColor(0x10,0xB9,0x81), "Watch live",  "All agents visible in one window, real-time status"),
    (RGBColor(0xF5,0x9E,0x0B), "Schedule",   "Set crons — agents work while you sleep"),
]

for i, (color, title, body) in enumerate(pills):
    x = Inches(0.5 + (i % 2) * 6.4)
    y = Inches(1.5 + (i // 2) * 2.7)
    r = rect(s, x, y, Inches(6.0), Inches(2.4), fill=CARD)
    bar = s.shapes.add_shape(1, x, y, Inches(0.18), Inches(2.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    txbox(s, title, x + Inches(0.4), y + Inches(0.2), Inches(5.4), Inches(0.55),
          size=22, bold=True, color=WHITE)
    txbox(s, body,  x + Inches(0.4), y + Inches(0.85), Inches(5.4), Inches(1.3),
          size=15, color=GREY)

slide_number(s, 3)

# ── Slide 4 — How It Works ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
accent_bar(s)
txbox(s, "⚙️  How It Works", Inches(0.6), Inches(0.3), Inches(10), Inches(0.9),
      size=40, bold=True, color=WHITE)

steps = [
    ("1", "Define Agent",    "CLAUDE.md with YAML frontmatter\n(name, emoji, args, MCPs, repos)"),
    ("2", "Spawn Instances", "Each instance clones repos into\nan isolated working directory"),
    ("3", "Run Tasks",       "Pass one-off args → Claude Code\nlaunches in a real PTY terminal"),
    ("4", "Watch & Resume",  "Live status: idle / thinking /\nwaiting. Resume from session ID"),
]

arrow_y = Inches(3.5)
for i, (num, title, body) in enumerate(steps):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.4)
    circle = s.shapes.add_shape(9, x, y, Inches(0.55), Inches(0.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
    txbox(s, num, x, y, Inches(0.55), Inches(0.55), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, title, x - Inches(0.1), y + Inches(0.65), Inches(2.8), Inches(0.5),
          size=17, bold=True, color=WHITE)
    txbox(s, body,  x - Inches(0.1), y + Inches(1.2), Inches(2.9), Inches(1.5),
          size=13, color=GREY)
    if i < 3:
        arr = s.shapes.add_shape(1, x + Inches(2.85), arrow_y - Inches(0.08),
                                  Inches(0.35), Inches(0.18))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT; arr.line.fill.background()

# Architecture box
rect(s, Inches(0.4), Inches(5.0), Inches(12.4), Inches(2.0), fill=CARD, line=ACCENT2)
txbox(s, "Electron + Svelte  ·  xterm.js + node-pty  ·  Supabase  ·  macOS launchd (crons)  ·  YAML agent definitions",
      Inches(0.7), Inches(5.1), Inches(12.0), Inches(0.6), size=15, color=ACCENT2, align=PP_ALIGN.CENTER)
txbox(s, "IPC-bridged architecture — all process/file/DB ops in main process, renderer is pure UI",
      Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.9), size=13, color=GREY, align=PP_ALIGN.CENTER)

slide_number(s, 4)

# ── Slide 5 — Key Features ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
accent_bar(s)
txbox(s, "✨  Key Features", Inches(0.6), Inches(0.3), Inches(10), Inches(0.9),
      size=40, bold=True, color=WHITE)

features = [
    ("🤖", "Specialised Agents",      "Custom instructions, whitelisted commands,\nrequired MCPs, auto-cloned repos per agent"),
    ("🖥️", "Live Multi-Terminal",     "xterm.js terminals — idle / thinking /\nwaiting status detected in real-time"),
    ("📅", "Cron Scheduling",          "Native macOS launchd plists — run agents\nautonomously on any schedule"),
    ("📜", "Session History",          "Every session logged (raw + plain-text).\nSearch, replay, and resume anytime"),
    ("🏪", "Agent Marketplace",        "Publish, star, fork, and share agent\ndefinitions across teams"),
    ("🔀", "Parallel Instances",       "Multiple isolated working directories\nper agent — true parallelism"),
]

for i, (emoji, title, body) in enumerate(features):
    col = i % 3
    row = i // 2
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.4 + row * 2.4)
    r = rect(s, x, y, Inches(4.05), Inches(2.2), fill=CARD, line=ACCENT)
    txbox(s, emoji, x + Inches(0.18), y + Inches(0.15), Inches(0.6), Inches(0.55), size=24)
    txbox(s, title, x + Inches(0.75), y + Inches(0.18), Inches(3.1), Inches(0.5),
          size=15, bold=True, color=WHITE)
    txbox(s, body,  x + Inches(0.18), y + Inches(0.75), Inches(3.7), Inches(1.3),
          size=12, color=GREY)

slide_number(s, 5)

# ── Slide 6 — Architecture Deep-Dive ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
accent_bar(s)
txbox(s, "🏗️  Architecture", Inches(0.6), Inches(0.3), Inches(10), Inches(0.9),
      size=40, bold=True, color=WHITE)

# Left column — layers
layers = [
    (ACCENT,  "Svelte UI (Renderer)",     "Multi-tab terminal, file browser, agent sidebar, modals"),
    (ACCENT2, "IPC Bridge (preload.ts)",  "Type-safe channel between renderer and main process"),
    (RGBColor(0x10,0xB9,0x81), "Electron Main",  "Agent Manager · PTY Manager · Session Manager\nCron Manager · MCP Manager · Marketplace Manager"),
    (RGBColor(0xF5,0x9E,0x0B), "Supabase",       "users · agents · instances · sessions · crons · marketplace"),
]
for i, (color, title, body) in enumerate(layers):
    y = Inches(1.4 + i * 1.45)
    r = rect(s, Inches(0.5), y, Inches(7.5), Inches(1.3), fill=CARD)
    bar2 = s.shapes.add_shape(1, Inches(0.5), y, Inches(0.18), Inches(1.3))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = color; bar2.line.fill.background()
    txbox(s, title, Inches(0.85), y + Inches(0.08), Inches(7.0), Inches(0.45),
          size=16, bold=True, color=WHITE)
    txbox(s, body,  Inches(0.85), y + Inches(0.55), Inches(7.0), Inches(0.65),
          size=12, color=GREY)

# Right column — file system
rect(s, Inches(8.5), Inches(1.4), Inches(4.4), Inches(5.7), fill=CARD, line=GREY)
txbox(s, "~/.swarm/", Inches(8.7), Inches(1.5), Inches(4.0), Inches(0.5),
      size=16, bold=True, color=ACCENT2)
fs_lines = [
    "CLAUDE.md          ← global rules",
    "agents/",
    "  <id>/",
    "    CLAUDE.md      ← agent def",
    "    instances/",
    "      0/           ← isolated dir",
    "        repo/      ← cloned repo",
    "        sessions/",
    "          *.log    ← raw output",
    "          *.txt    ← plain text",
]
txbox(s, "\n".join(fs_lines), Inches(8.7), Inches(2.1), Inches(4.0), Inches(4.8),
      size=11, color=GREY)

slide_number(s, 6)

# ── Slide 7 — Traction / Demo ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
accent_bar(s)
txbox(s, "🚀  What We Built in the Hackathon", Inches(0.6), Inches(0.3), Inches(12), Inches(0.9),
      size=36, bold=True, color=WHITE)

bullets = [
    ("✅", "Full Electron desktop app — installable, distributable"),
    ("✅", "Agent CRUD with YAML frontmatter definitions"),
    ("✅", "Parallel PTY terminals with live idle/thinking/waiting status"),
    ("✅", "Session logging, search, and resume (--resume flag)"),
    ("✅", "macOS cron scheduling via launchd plist generation"),
    ("✅", "MCP server connectivity checks before session launch"),
    ("✅", "Community marketplace — publish, fork, star agents"),
    ("✅", "File browser with content preview inside instances"),
]

for i, (icon, text) in enumerate(bullets):
    y = Inches(1.4 + i * 0.7)
    txbox(s, icon, Inches(0.6), y, Inches(0.5), Inches(0.6), size=20, color=RGBColor(0x10,0xB9,0x81))
    txbox(s, text, Inches(1.2), y, Inches(11.5), Inches(0.6), size=18, color=WHITE)

slide_number(s, 7)

# ── Slide 8 — What's Next ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
accent_bar(s)
txbox(s, "🔭  What's Next", Inches(0.6), Inches(0.3), Inches(10), Inches(0.9),
      size=40, bold=True, color=WHITE)

nexts = [
    ("🔗", "Agent Chaining",     "Output of one agent triggers the next — automated pipelines"),
    ("📊", "Dashboard",          "Cross-agent analytics: cost, success rate, run time"),
    ("🌐", "Cloud Execution",    "Run agents on remote sandboxes, not just local machine"),
    ("🔔", "Notifications",      "Slack / email alerts when an agent finishes or fails"),
    ("🔑", "Team Workspaces",    "Shared agent libraries with role-based access"),
    ("📦", "One-click Deploy",   "Export an agent + its environment as a shareable bundle"),
]

for i, (emoji, title, body) in enumerate(nexts):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.4 + row * 2.6)
    r = rect(s, x, y, Inches(4.05), Inches(2.3), fill=CARD, line=ACCENT2)
    txbox(s, emoji, x + Inches(0.2), y + Inches(0.15), Inches(0.65), Inches(0.65), size=26)
    txbox(s, title, x + Inches(0.85), y + Inches(0.2), Inches(3.0), Inches(0.5),
          size=16, bold=True, color=WHITE)
    txbox(s, body,  x + Inches(0.2), y + Inches(0.85), Inches(3.7), Inches(1.2),
          size=13, color=GREY)

slide_number(s, 8)

# ── Slide 9 — Closing ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)

blob2 = s.shapes.add_shape(9, Inches(2), Inches(1), Inches(9), Inches(6))
blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(0x1A,0x00,0x50)
blob2.line.fill.background()

txbox(s, "🐝", Inches(5.4), Inches(1.5), Inches(2.5), Inches(1.5), size=80, align=PP_ALIGN.CENTER)
txbox(s, "SWARM", Inches(1), Inches(2.9), W - Inches(2), Inches(1.5),
      size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Your agents. Your rules. Running in parallel.", Inches(1), Inches(4.4), W - Inches(2), Inches(0.8),
      size=22, color=ACCENT2, align=PP_ALIGN.CENTER)
txbox(s, "github.com/Nikhil-Chavanke-21/swarm", Inches(1), Inches(5.3), W - Inches(2), Inches(0.6),
      size=16, color=GREY, align=PP_ALIGN.CENTER)
accent_bar(s)

OUT = "/Users/nc/regie/swarm/Swarm_Hackathon_Deck.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
